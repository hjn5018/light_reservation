import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(filename="presentation.pptx"):
    prs = Presentation()
    
    # 16:9 Widescreen slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Color Palette (Dark Theme)
    BG_COLOR = RGBColor(15, 23, 42)        # #0F172A Slate 900
    PANEL_COLOR = RGBColor(30, 41, 59)     # #1E293B Slate 800
    TEXT_WHITE = RGBColor(248, 250, 252)   # #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184)   # #94A3B8 Cool Gray
    COLOR_EMERALD = RGBColor(16, 185, 129) # #10B981 Emerald 500
    COLOR_AMBER = RGBColor(245, 158, 11)   # #F59E0B Amber 500
    COLOR_BLUE = RGBColor(59, 130, 246)    # #3B82F6 Blue 500
    COLOR_RED = RGBColor(239, 68, 68)      # #EF4444 Red 500
    
    FONT_TITLE = "Malgun Gothic"
    FONT_BODY = "Malgun Gothic"

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_slide_header(slide, title_text, category_text="Smart Multi-Resource Reservation IoT System"):
        # Top-left small category tracker
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = FONT_BODY
        p_cat.font.size = Pt(10)
        p_cat.font.color.rgb = COLOR_EMERALD
        p_cat.font.bold = True
        
        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(10), Inches(0.7))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(24)
        p_title.font.color.rgb = TEXT_WHITE
        p_title.font.bold = True
        
        # Bottom footer
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.3))
        tf_foot = footer_box.text_frame
        tf_foot.margin_left = tf_foot.margin_right = tf_foot.margin_top = tf_foot.margin_bottom = 0
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "IoT 제어 실습  |  2-Pi 스마트 자원 예약 시스템 발표자료"
        p_foot.font.name = FONT_BODY
        p_foot.font.size = Pt(9)
        p_foot.font.color.rgb = TEXT_MUTED

    # Helpers to safely add images
    def try_add_image(slide, path, left, top, width, height):
        # Resolve path relative to this script's directory
        resolved_path = path
        if not os.path.isabs(path):
            base_dir = os.path.dirname(os.path.abspath(__file__))
            # Path might be "../images/..." or "images/..."
            # Strip leading ../ or similar
            clean_path = path.lstrip("./").replace("../", "")
            resolved_path = os.path.join(os.path.dirname(base_dir), clean_path)
            
        if os.path.exists(resolved_path):
            try:
                slide.shapes.add_picture(resolved_path, left, top, width, height)
                return True
            except Exception as e:
                print(f"Error adding image {resolved_path}: {e}")
        else:
            print(f"Warning: Image not found at {resolved_path}")
            # Draw a placeholder box
            ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            ph.fill.solid()
            ph.fill.fore_color.rgb = PANEL_COLOR
            ph.line.color.rgb = TEXT_MUTED
            tf = ph.text_frame
            p = tf.paragraphs[0]
            p.text = f"Image Placeholder:\n{os.path.basename(path)}"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_WHITE
        return False

    slide_layout = prs.slide_layouts[6]  # Blank Layout

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Cover)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide1)
    
    # Accent Bar
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.1), Inches(3.5), Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_EMERALD
    accent_bar.line.fill.background()
    
    # Main Title Text
    title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(2.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "IoT 기반 스마트 통합 자원 예약 시스템"
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(38)
    p1.font.color.rgb = TEXT_WHITE
    p1.font.bold = True
    
    p2 = tf1.add_paragraph()
    p2.text = "Smart Multi-Resource Reservation IoT System"
    p2.font.name = FONT_TITLE
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_AMBER
    p2.font.bold = True
    p2.space_before = Pt(6)

    # Subtitle / Meta info
    sub_box = slide1.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.7), Inches(1.5))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub1 = tf_sub.paragraphs[0]
    p_sub1.text = "2대 라즈베리 파이 네트워크 연동 및 3대 하드웨어 피드백 기반 공용 자원 관리 솔루션"
    p_sub1.font.name = FONT_BODY
    p_sub1.font.size = Pt(14)
    p_sub1.font.color.rgb = TEXT_MUTED
    
    p_sub2 = tf_sub.add_paragraph()
    p_sub2.text = "과목: IoT 제어 실습  |  발표자: 홍길동 (편집 가능)"
    p_sub2.font.name = FONT_BODY
    p_sub2.font.size = Pt(13)
    p_sub2.font.color.rgb = COLOR_EMERALD
    p_sub2.space_before = Pt(20)

    # ----------------------------------------------------
    # SLIDE 2: Table of Contents (목차)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide2)
    add_slide_header(slide2, "발표 목차 (Contents)")
    
    start_top = Inches(1.6)
    panel_h = Inches(0.85)
    gap = Inches(0.2)
    
    toc_items = [
        ("01", "📋 프로젝트 개요 및 배경", "공용 자원 무단 점유와 예약 No-Show 분쟁의 원인과 시스템적 대안"),
        ("02", "🛠️ 사용 장비 및 시스템 구성", "2대 라즈베리 파이 및 3대 물리 액추에이터 하드웨어/네트워크 연동"),
        ("03", "📅 개발 일정 및 전체 흐름도", "3주 개발 주기 타임라인 및 클라이언트-서버 간 제어 순서도"),
        ("04", "🔒 핵심 소프트웨어 기술 설명", "CGI 게이트웨이, TCP 소켓 통신, 멀티스레딩, Mutex 데이터 잠금"),
        ("05", "📷 작동 시연 및 기대효과", "실제 시스템 구동 화면 구성, 기대효과 및 향후 발전 가능성")
    ]
    
    for i, (num, title, desc) in enumerate(toc_items):
        t_top = start_top + i * (panel_h + gap)
        
        # Round panel
        panel = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), t_top, Inches(11.733), panel_h)
        panel.fill.solid()
        panel.fill.fore_color.rgb = PANEL_COLOR
        panel.line.color.rgb = COLOR_EMERALD if i%2==0 else COLOR_AMBER
        panel.line.width = Pt(1.5)
        
        tf_panel = panel.text_frame
        tf_panel.word_wrap = True
        tf_panel.margin_left = Inches(0.3)
        tf_panel.margin_top = Inches(0.12)
        
        p = tf_panel.paragraphs[0]
        # Slide number badge
        r_num = p.add_run()
        r_num.text = f"{num}  "
        r_num.font.name = FONT_TITLE
        r_num.font.size = Pt(18)
        r_num.font.bold = True
        r_num.font.color.rgb = COLOR_EMERALD if i%2==0 else COLOR_AMBER
        
        # Title
        r_title = p.add_run()
        r_title.text = f"{title}   |   "
        r_title.font.name = FONT_TITLE
        r_title.font.size = Pt(16)
        r_title.font.bold = True
        r_title.font.color.rgb = TEXT_WHITE
        
        # Description
        r_desc = p.add_run()
        r_desc.text = desc
        r_desc.font.name = FONT_BODY
        r_desc.font.size = Pt(12)
        r_desc.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 3: Development Background & Problems
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide3)
    add_slide_header(slide3, "1. 개발 배경 및 필요성")
    
    callout = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(6.8), Inches(0.7))
    callout.fill.solid()
    callout.fill.fore_color.rgb = PANEL_COLOR
    callout.line.color.rgb = COLOR_AMBER
    tf_call = callout.text_frame
    tf_call.word_wrap = True
    tf_call.margin_left = Inches(0.2)
    p_call = tf_call.paragraphs[0]
    p_call.text = "💡 공용 자원 관리의 고질적인 무단 점유와 예약 방치(No-Show)를 시스템 제어로 해결"
    p_call.font.name = FONT_TITLE
    p_call.font.size = Pt(11.5)
    p_call.font.color.rgb = TEXT_WHITE
    p_call.font.bold = True
    
    # 2 columns structure
    left_c = slide3.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(6.8), Inches(4.5))
    tf_l = left_c.text_frame
    tf_l.word_wrap = True
    
    p_lh = tf_l.paragraphs[0]
    p_lh.text = "🚨 기존 공용 자원의 관리 한계"
    p_lh.font.name = FONT_TITLE
    p_lh.font.size = Pt(15)
    p_lh.font.color.rgb = COLOR_AMBER
    p_lh.font.bold = True
    p_lh.space_after = Pt(8)
    
    probs = [
        "비인가 무단 점유: 실시간 예약 및 점유 상태가 불투명하여 발생하는 자원 가로채기와 마찰",
        "자원 독점과 노쇼: 예약 후 방치해 두어 다른 사용자들의 자원 접근 기회를 일방적으로 박탈",
        "물리 피드백 부재: 화면을 모니터링하기 전까지는 기기 상태나 반납 시간을 직관적으로 알기 어려움"
    ]
    for text in probs:
        p = tf_l.add_paragraph()
        p.text = "• " + text
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(6)
        
    p_rh = tf_l.add_paragraph()
    p_rh.text = "✅ 스마트 예약 시스템의 해결 대안"
    p_rh.font.name = FONT_TITLE
    p_rh.font.size = Pt(15)
    p_rh.font.color.rgb = COLOR_EMERALD
    p_rh.font.bold = True
    p_rh.space_before = Pt(14)
    p_rh.space_after = Pt(8)
    
    sols = [
        "웹 실시간 대시보드: 5초 주기로 모든 공용 물품/공간의 예약 상태를 한눈에 비동기 모니터링",
        "4자리 PIN 본인 식별: 복잡한 가입 없이 예약 시 PIN 번호를 활용해 타인의 오용 및 해제 원천 차단",
        "자동 만료 스케줄러: 지정 대기 시간 내 미사용 시 노쇼 취소 처리, 최대 사용시간 경과 시 자동 반납"
    ]
    for text in sols:
        p = tf_l.add_paragraph()
        p.text = "• " + text
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(6)
        
    # AI Illustration on the Right
    try_add_image(slide3, "../images/background_need_illustration.png", Inches(8.0), Inches(1.5), Inches(4.5), Inches(4.5))

    # ----------------------------------------------------
    # SLIDE 4: Equipment Used (사용 장비 및 부품)
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide4)
    add_slide_header(slide4, "2. 사용 장비 및 부품 사양")
    
    grid_w = Inches(2.7)
    grid_h = Inches(4.5)
    grid_top = Inches(1.8)
    grid_gap = Inches(0.3)
    
    equipments = [
        ("라즈베리 파이 4 (2대)", "시스템 연동 제어 코어", 
         ["Pi A: Apache 웹 서버 호스팅 및 웹-서버 CGI 소켓 게이트웨이 역할", 
          "Pi B: Flask REST API 처리, 백그라운드 소켓 수신 및 하드웨어 액추에이터 제어"],
         "../images/pi_A_pi_B_view.jpg"),
        ("삼색 LED (3개)", "직관적인 시각 피드백", 
         ["자원의 상태 변화를 신호등 개념으로 가시화", 
          "초록 LED: 대기 상태(Available)", 
          "노란 LED: 예약 상태(Reserved)", 
          "빨간 LED: 사용 상태(In Use)"],
         "../images/breadboard_view.jpg"),
        ("16x2 Character LCD", "포커스 자원 문자 출력", 
         ["선택된 기기의 영문 자원명 실시간 렌더링", 
          "기기의 현재 세부 동작 상태 문자 표현", 
          "가변저항 회로 연동을 통해 전압 대비 미세 튜닝"],
         "../images/pi_B_view.jpg"),
        ("피에조 패시브 부저", "상태 전이 청각 피드백", 
         ["주파수 PWM 변조를 활용한 상태 변경 음계 연주", 
          "사용 시작 시: 상승음 (도-미-솔) 멜로디 알림", 
          "자동 만료/반납 시: 하강음 (솔-미-도) 멜로디 경보"],
         "../images/breadboard_view.jpg")
    ]
    
    for idx, (name, caption, specs, img_path) in enumerate(equipments):
        left_pos = Inches(0.8) + idx * (grid_w + grid_gap)
        
        # Grid box
        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, grid_top, grid_w, grid_h)
        box.fill.solid()
        box.fill.fore_color.rgb = PANEL_COLOR
        box.line.color.rgb = COLOR_EMERALD if idx%2==0 else COLOR_AMBER
        box.line.width = Pt(1.5)
        
        # Add Image inside card
        img_h = Inches(1.3)
        try_add_image(slide4, img_path, left_pos + Inches(0.15), grid_top + Inches(0.15), grid_w - Inches(0.3), img_h)
        
        # Text box below image
        tb = slide4.shapes.add_textbox(left_pos + Inches(0.1), grid_top + img_h + Inches(0.2), grid_w - Inches(0.2), grid_h - img_h - Inches(0.3))
        tf_box = tb.text_frame
        tf_box.word_wrap = True
        tf_box.margin_left = tf_box.margin_right = tf_box.margin_top = tf_box.margin_bottom = 0
        
        # Name
        p_name = tf_box.paragraphs[0]
        p_name.text = name
        p_name.alignment = PP_ALIGN.CENTER
        p_name.font.name = FONT_TITLE
        p_name.font.size = Pt(12)
        p_name.font.bold = True
        p_name.font.color.rgb = TEXT_WHITE
        p_name.space_after = Pt(2)
        
        # Caption
        p_cap = tf_box.add_paragraph()
        p_cap.text = caption
        p_cap.alignment = PP_ALIGN.CENTER
        p_cap.font.name = FONT_BODY
        p_cap.font.size = Pt(9.5)
        p_cap.font.color.rgb = COLOR_AMBER if idx%2==1 else COLOR_EMERALD
        p_cap.space_after = Pt(8)
        
        # Bullet Specs
        for spec in specs:
            p_spec = tf_box.add_paragraph()
            p_spec.text = "• " + spec
            p_spec.font.name = FONT_BODY
            p_spec.font.size = Pt(8.5)
            p_spec.font.color.rgb = TEXT_WHITE
            p_spec.space_after = Pt(3)

    # ----------------------------------------------------
    # SLIDE 5: Development Schedule (Gantt Chart)
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide5)
    add_slide_header(slide5, "3. 개발 일정 및 수행 내역 (Gantt Chart)")
    
    # Custom Gantt Chart representation
    # Table headers
    hdr_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.4))
    tf_hdr = hdr_box.text_frame
    tf_hdr.margin_left = tf_hdr.margin_right = tf_hdr.margin_top = tf_hdr.margin_bottom = 0
    p_hdr = tf_hdr.paragraphs[0]
    
    # Grid math
    gantt_left = Inches(4.2)
    gantt_width = Inches(8.333)
    w_width = gantt_width / 3.0  # 2.777 inches per week
    
    # Draw headers for weeks
    weeks = [
        ("1주차 (06.04 ~ 06.10)", gantt_left),
        ("2주차 (06.11 ~ 06.17)", gantt_left + w_width),
        ("3주차 (06.18 ~ 06.25)", gantt_left + w_width * 2)
    ]
    for w_name, w_pos in weeks:
        wb = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, w_pos + Inches(0.05), Inches(1.5), w_width - Inches(0.1), Inches(0.4))
        wb.fill.solid()
        wb.fill.fore_color.rgb = PANEL_COLOR
        wb.line.color.rgb = COLOR_EMERALD
        tf_wb = wb.text_frame
        p_wb = tf_wb.paragraphs[0]
        p_wb.text = w_name
        p_wb.font.size = Pt(10)
        p_wb.font.bold = True
        p_wb.font.color.rgb = TEXT_WHITE
        p_wb.alignment = PP_ALIGN.CENTER
        
    # Gantt rows
    tasks = [
        ("📋 요구사항 정의 및 기획", 1, COLOR_EMERALD, "자원 기획 및 명세 정의"),
        ("⚡ 회로 설계 및 부품 구성", 1, COLOR_EMERALD, "GPIO 회로 배선도 수립"),
        ("📟 Web & CGI 서버 개발 (Pi A)", 2, COLOR_AMBER, "Apache 웹 리소스 & CGI 개발"),
        ("⚙️ Flask & Socket 개발 (Pi B)", 2, COLOR_AMBER, "REST API 및 소켓 리스너 개발"),
        ("🚨 하드웨어 제어 모듈 코딩", 2, COLOR_AMBER, "LED/LCD/Buzzer 드라이버 구현"),
        ("🧪 통합 연동 및 디버깅", 3, COLOR_BLUE, "소켓 통신 데이터 검증"),
        ("🔒 스케줄러 안정화 및 배포", 3, COLOR_BLUE, "데몬 서비스 등록 & setup.sh 제작")
    ]
    
    start_top = Inches(2.1)
    row_height = Inches(0.55)
    
    for idx, (t_name, week_idx, color, desc) in enumerate(tasks):
        top_pos = start_top + idx * row_height
        
        # Row background panel for task label
        lbl_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, Inches(3.2), Inches(0.45))
        lbl_box.fill.solid()
        lbl_box.fill.fore_color.rgb = PANEL_COLOR
        lbl_box.line.color.rgb = TEXT_MUTED
        lbl_box.line.width = Pt(1.0)
        tf_lbl = lbl_box.text_frame
        p_lbl = tf_lbl.paragraphs[0]
        p_lbl.text = t_name
        p_lbl.font.size = Pt(10)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = TEXT_WHITE
        p_lbl.alignment = PP_ALIGN.LEFT
        
        # Gantt Bar
        bar_left = gantt_left + (week_idx - 1) * w_width
        bar_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bar_left + Inches(0.05), top_pos, w_width - Inches(0.1), Inches(0.45))
        bar_box.fill.solid()
        bar_box.fill.fore_color.rgb = color
        bar_box.line.fill.background()
        tf_bar = bar_box.text_frame
        p_bar = tf_bar.paragraphs[0]
        p_bar.text = desc
        p_bar.font.size = Pt(9.5)
        p_bar.font.bold = True
        p_bar.font.color.rgb = BG_COLOR
        p_bar.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 6: Overall System Structure (UML Component Diagram)
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide6)
    add_slide_header(slide6, "4. 전체 시스템 구성도 (UML Component Diagram)")
    
    # Left component: Pi A
    pi_a = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(3.8), Inches(4.3))
    pi_a.fill.solid()
    pi_a.fill.fore_color.rgb = PANEL_COLOR
    pi_a.line.color.rgb = COLOR_EMERALD
    pi_a.line.width = Pt(2.0)
    tf_a = pi_a.text_frame
    tf_a.word_wrap = True
    p_a = tf_a.paragraphs[0]
    p_a.text = "«component»\nPi A (Web Gateway)"
    p_a.alignment = PP_ALIGN.CENTER
    p_a.font.name = FONT_TITLE
    p_a.font.size = Pt(14)
    p_a.font.bold = True
    p_a.font.color.rgb = COLOR_EMERALD
    
    # Pi A Sub-components
    a_subs = [
        "«component» Dashboard UI\n(index.html & app.js)",
        "«component» CGI Client\n(update_status.py)"
    ]
    for s_idx, text in enumerate(a_subs):
        sb = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.1), Inches(2.8) + s_idx * Inches(1.2), Inches(3.2), Inches(0.9))
        sb.fill.solid()
        sb.fill.fore_color.rgb = BG_COLOR
        sb.line.color.rgb = TEXT_MUTED
        sb.line.width = Pt(1.0)
        tf_sb = sb.text_frame
        p_sb = tf_sb.paragraphs[0]
        p_sb.text = text
        p_sb.font.size = Pt(9.5)
        p_sb.font.color.rgb = TEXT_WHITE
        p_sb.alignment = PP_ALIGN.CENTER
        
    # Right component: Pi B
    pi_b = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.7), Inches(1.8), Inches(3.8), Inches(4.3))
    pi_b.fill.solid()
    pi_b.fill.fore_color.rgb = PANEL_COLOR
    pi_b.line.color.rgb = COLOR_AMBER
    pi_b.line.width = Pt(2.0)
    tf_b = pi_b.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    p_b.text = "«component»\nPi B (Hardware Core)"
    p_b.alignment = PP_ALIGN.CENTER
    p_b.font.name = FONT_TITLE
    p_b.font.size = Pt(14)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_AMBER
    
    # Pi B Sub-components
    b_subs = [
        "«component» Flask REST API",
        "«component» Socket Listener (Thread)",
        "«component» Mutex Lock (state_lock) & data.json"
    ]
    for s_idx, text in enumerate(b_subs):
        sb = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(2.6) + s_idx * Inches(0.85), Inches(3.2), Inches(0.65))
        sb.fill.solid()
        sb.fill.fore_color.rgb = BG_COLOR
        sb.line.color.rgb = TEXT_MUTED
        sb.line.width = Pt(1.0)
        tf_sb = sb.text_frame
        p_sb = tf_sb.paragraphs[0]
        p_sb.text = text
        p_sb.font.size = Pt(9.0)
        p_sb.font.color.rgb = TEXT_WHITE
        p_sb.alignment = PP_ALIGN.CENTER
        
    # Pi B Port section
    pb_port = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(5.2), Inches(3.2), Inches(0.5))
    pb_port.fill.solid()
    pb_port.fill.fore_color.rgb = BG_COLOR
    pb_port.line.color.rgb = COLOR_AMBER
    pb_port.line.width = Pt(1.0)
    tf_port = pb_port.text_frame
    p_port = tf_port.paragraphs[0]
    p_port.text = "«port» GPIO ➔ LEDs/LCD/Buzzer"
    p_port.font.size = Pt(9.0)
    p_port.font.bold = True
    p_port.font.color.rgb = COLOR_AMBER
    p_port.alignment = PP_ALIGN.CENTER
    
    # Middle arrows/labels
    # TCP command
    tcp_arrow = slide6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.8), Inches(3.0), Inches(3.7), Inches(0.3))
    tcp_arrow.fill.solid()
    tcp_arrow.fill.fore_color.rgb = COLOR_EMERALD
    tcp_arrow.line.fill.background()
    
    tcp_txt = slide6.shapes.add_textbox(Inches(4.8), Inches(2.5), Inches(3.7), Inches(0.4))
    tf_tcp = tcp_txt.text_frame
    p_tcp = tf_tcp.paragraphs[0]
    p_tcp.text = "«interface» TCP:50007 (Command)"
    p_tcp.font.size = Pt(10)
    p_tcp.font.bold = True
    p_tcp.font.color.rgb = TEXT_WHITE
    p_tcp.alignment = PP_ALIGN.CENTER
    
    # HTTP polling
    http_arrow = slide6.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(4.8), Inches(4.5), Inches(3.7), Inches(0.3))
    http_arrow.fill.solid()
    http_arrow.fill.fore_color.rgb = COLOR_AMBER
    http_arrow.line.fill.background()
    
    http_txt = slide6.shapes.add_textbox(Inches(4.8), Inches(4.0), Inches(3.7), Inches(0.4))
    tf_http = http_txt.text_frame
    p_http = tf_http.paragraphs[0]
    p_http.text = "«interface» HTTP:5000 (Polling)"
    p_http.font.size = Pt(10)
    p_http.font.bold = True
    p_http.font.color.rgb = TEXT_WHITE
    p_http.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 7: Overall System Flowchart (UML Sequence Diagram)
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide7)
    add_slide_header(slide7, "5. 전체 시스템 연동 흐름도 (UML Sequence Diagram)")
    
    # X coordinates of 4 lifelines
    # Col 1 (User): 2.26, Col 2 (Pi A): 5.20, Col 3 (Pi B): 8.13, Col 4 (Hardware): 11.06
    cols_x = [Inches(2.26), Inches(5.20), Inches(8.13), Inches(11.06)]
    names = [
        "🧑‍💻 User\n(Dashboard)",
        "📟 Pi A\n(Gateway)",
        "⚙️ Pi B\n(Core Server)",
        "🚨 Hardware\n(Actuators)"
    ]
    
    # Draw lifeline headers
    for idx, x_pos in enumerate(cols_x):
        h_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos - Inches(0.9), Inches(1.5), Inches(1.8), Inches(0.6))
        h_box.fill.solid()
        h_box.fill.fore_color.rgb = PANEL_COLOR
        h_box.line.color.rgb = COLOR_EMERALD if idx%2==0 else COLOR_AMBER
        h_box.line.width = Pt(1.5)
        tf_h = h_box.text_frame
        p_h = tf_h.paragraphs[0]
        p_h.text = names[idx]
        p_h.font.size = Pt(10)
        p_h.font.bold = True
        p_h.font.color.rgb = TEXT_WHITE
        p_h.alignment = PP_ALIGN.CENTER
        
        # Vertical lifeline line
        line = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos - Inches(0.015), Inches(2.1), Inches(0.03), Inches(4.3))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(80, 80, 80)
        line.line.fill.background()
        
    # Message definitions
    # (From Col, To Col, Text, Style [Solid/Dashed], Offset Top)
    messages = [
        (1, 2, "1. 예약/사용 요청 (PIN 입력)", "solid", Inches(2.3)),
        (2, 3, "2. TCP Socket 전송 (Port 50007)", "solid", Inches(2.8)),
        (3, 3, "3. Mutex Lock 획득 & DB 갱신", "self", Inches(3.3)),
        (3, 4, "4. GPIO 제어 신호 출력", "solid", Inches(3.8)),
        (4, 1, "5. LED 점등, 부저 멜로디 및 LCD 텍스트 피드백", "dashed", Inches(4.3)),
        (3, 3, "6. 만료 스케줄러 (No-show/시간초과 회수)", "self", Inches(4.8)),
        (1, 3, "7. HTTP REST API 조회 (5초 주기)", "solid", Inches(5.3)),
        (3, 1, "8. JSON 전역 상태 & 카운트다운 정보 반환", "dashed", Inches(5.8))
    ]
    
    for from_col, to_col, msg_txt, style, m_top in messages:
        x_from = cols_x[from_col - 1]
        x_to = cols_x[to_col - 1]
        
        # Self call
        if style == "self":
            # Right loop
            loop = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_from, m_top, Inches(0.4), Inches(0.35))
            loop.fill.solid()
            loop.fill.fore_color.rgb = PANEL_COLOR
            loop.line.color.rgb = COLOR_AMBER
            loop.line.width = Pt(1.2)
            
            # Label
            lbl = slide7.shapes.add_textbox(x_from + Inches(0.45), m_top - Inches(0.05), Inches(3.0), Inches(0.4))
            tf_l = lbl.text_frame
            p_l = tf_l.paragraphs[0]
            p_l.text = msg_txt
            p_l.font.size = Pt(8.5)
            p_l.font.bold = True
            p_l.font.color.rgb = COLOR_AMBER
        else:
            # Line
            line_l = min(x_from, x_to)
            line_w = abs(x_from - x_to)
            
            arrow_shape = MSO_SHAPE.RIGHT_ARROW if x_to > x_from else MSO_SHAPE.LEFT_ARROW
            arrow = slide7.shapes.add_shape(arrow_shape, line_l, m_top, line_w, Inches(0.12))
            arrow.fill.solid()
            # Dashed uses blue, Solid uses emerald
            arrow.fill.fore_color.rgb = COLOR_BLUE if style == "dashed" else COLOR_EMERALD
            arrow.line.fill.background()
            
            # Text above line
            tb = slide7.shapes.add_textbox(line_l, m_top - Inches(0.32), line_w, Inches(0.35))
            tf_b = tb.text_frame
            tf_b.word_wrap = True
            tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = 0
            p_b = tf_b.paragraphs[0]
            p_b.text = msg_txt
            p_b.font.size = Pt(8.5)
            p_b.font.bold = True
            p_b.font.color.rgb = TEXT_WHITE
            p_b.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 8: Detailed Technology Explanations (CGI, Socket, Thread, Mutex)
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide8)
    add_slide_header(slide8, "6. 주요 기술 명세 및 동기화 구현")
    
    techs = [
        ("📡 CGI 연동 기술", "Apache & ExecCGI", 
         ["• 웹 클라이언트(JS)와 백엔드 서버 간의 데이터 게이트웨이 스크립트 동작",
          "• 웹 요청 환경변수를 파싱해 서버 소켓 연결로 이어주는 동적 커넥터 역할 수행"]),
        ("🔌 저지연 TCP 소켓", "포트 50007 제어 릴레이", 
         ["• HTTP 프로토콜의 헤더 오버헤드를 생략하고 원시 TCP 커넥션을 맺어 제어 지연시간 최소화",
          "• Pi A CGI에서 Pi B 소켓 리스너로 실시간 JSON 문자열 패킷 즉각 송신"]),
        ("🧵 멀티스레딩", "Daemon Threads 구현", 
         ["• 소켓 리스너 스레드와 타이머 스케줄러 스레드가 메인 API(Flask)와 별개로 병렬 실행",
          "• 데몬 모드로 작동하여 메인 시스템 종료 시 부작용 없이 동시 회수 및 종료"]),
        ("🔒 뮤텍스 락 (Mutex Lock)", "상호 배제 (Lock)", 
         ["• threading.Lock() 객체를 생성하여 데이터 및 물리 기기 동시 제어 충돌 차단",
          "• 소켓 쓰기 동작 중 불완전한 상태 정보를 Flask API가 읽지 못하게 크리티컬 섹션 완벽 보호"])
    ]
    
    grid_w = Inches(2.7)
    grid_h = Inches(4.5)
    grid_top = Inches(1.8)
    grid_gap = Inches(0.3)
    
    for idx, (title, caption, specs) in enumerate(techs):
        left_pos = Inches(0.8) + idx * (grid_w + grid_gap)
        
        # Grid box
        box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, grid_top, grid_w, grid_h)
        box.fill.solid()
        box.fill.fore_color.rgb = PANEL_COLOR
        box.line.color.rgb = COLOR_EMERALD if idx%2==0 else COLOR_AMBER
        box.line.width = Pt(1.5)
        
        tf_box = box.text_frame
        tf_box.word_wrap = True
        tf_box.margin_left = Inches(0.2)
        tf_box.margin_right = Inches(0.2)
        tf_box.margin_top = Inches(0.2)
        
        # Title
        p_name = tf_box.paragraphs[0]
        p_name.text = title
        p_name.font.name = FONT_TITLE
        p_name.font.size = Pt(13)
        p_name.font.bold = True
        p_name.font.color.rgb = COLOR_EMERALD if idx%2==0 else COLOR_AMBER
        p_name.space_after = Pt(4)
        
        # Caption
        p_cap = tf_box.add_paragraph()
        p_cap.text = caption
        p_cap.font.name = FONT_BODY
        p_cap.font.size = Pt(11)
        p_cap.font.color.rgb = TEXT_MUTED
        p_cap.space_after = Pt(12)
        
        # Bullet Specs
        for spec in specs:
            p_spec = tf_box.add_paragraph()
            p_spec.text = spec
            p_spec.font.name = FONT_BODY
            p_spec.font.size = Pt(9.5)
            p_spec.font.color.rgb = TEXT_WHITE
            p_spec.space_after = Pt(6)

    # ----------------------------------------------------
    # SLIDE 9: Actual Operating Image (실제 작동 이미지 그리드)
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide9)
    add_slide_header(slide9, "7. 실제 시스템 작동 시연 화면")
    
    # Left Col (Width: 5.6)
    try_add_image(slide9, "../images/full_system_view.jpg", Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.1))
    lbl1 = slide9.shapes.add_textbox(Inches(0.8), Inches(3.9), Inches(5.6), Inches(0.3))
    lbl1.text_frame.paragraphs[0].text = "🖥️ 전체 연동 시스템 (라즈베리파이 2대 + 브레드보드 + 대시보드 브라우저)"
    lbl1.text_frame.paragraphs[0].font.size = Pt(8.5)
    lbl1.text_frame.paragraphs[0].font.color.rgb = TEXT_MUTED
    
    try_add_image(slide9, "../images/browser_pi_A_pi_B_view.png", Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.1))
    lbl2 = slide9.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(5.6), Inches(0.3))
    lbl2.text_frame.paragraphs[0].text = "💻 웹 대시보드 및 터미널 모니터링 (Apache 웹 브라우저 + SSH 터미널)"
    lbl2.text_frame.paragraphs[0].font.size = Pt(8.5)
    lbl2.text_frame.paragraphs[0].font.color.rgb = TEXT_MUTED

    # Right Col (Width: 5.6)
    try_add_image(slide9, "../images/pi_A_pi_B_view.jpg", Inches(6.9), Inches(1.8), Inches(5.6), Inches(2.1))
    lbl3 = slide9.shapes.add_textbox(Inches(6.9), Inches(3.9), Inches(5.6), Inches(0.3))
    lbl3.text_frame.paragraphs[0].text = "🔌 라즈베리파이 2대 및 회로 구성 (네트워크 연동 및 브레드보드 연결)"
    lbl3.text_frame.paragraphs[0].font.size = Pt(8.5)
    lbl3.text_frame.paragraphs[0].font.color.rgb = TEXT_MUTED
    
    try_add_image(slide9, "../images/pi_B_view.jpg", Inches(6.9), Inches(4.3), Inches(2.7), Inches(2.1))
    lbl4 = slide9.shapes.add_textbox(Inches(6.9), Inches(6.4), Inches(2.7), Inches(0.3))
    lbl4.text_frame.paragraphs[0].text = "⚙️ Pi B 제어부"
    lbl4.text_frame.paragraphs[0].font.size = Pt(8.5)
    lbl4.text_frame.paragraphs[0].font.color.rgb = TEXT_MUTED
    
    try_add_image(slide9, "../images/breadboard_view.jpg", Inches(9.8), Inches(4.3), Inches(2.7), Inches(2.1))
    lbl5 = slide9.shapes.add_textbox(Inches(9.8), Inches(6.4), Inches(2.7), Inches(0.3))
    lbl5.text_frame.paragraphs[0].text = "🚨 액추에이터 회로"
    lbl5.text_frame.paragraphs[0].font.size = Pt(8.5)
    lbl5.text_frame.paragraphs[0].font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 10: Future Potential (향후 발전 가능성)
    # ----------------------------------------------------
    slide10 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide10)
    add_slide_header(slide10, "8. 향후 시스템 발전 가능성")
    
    potentials = [
        ("🌐 클라우드 데이터베이스 연동 확장", "로컬 data.json 파일 DB를 클라우드 RDB/NoSQL 구조로 전면 업그레이드. 다수의 보관함 및 지점 자원을 글로벌 단위 단일 망 통합 모니터링 구조로 확장 가능."),
        ("📱 전용 모바일 앱 및 하이브리드 PWA 연동", "PWA 및 모바일 전용 하이브리드 앱을 배포하여 접근성 극대화. 사용 기기 반납 1분 전 리마인드 푸시 및 비상시 강제 만료 즉각 앱 알림 송출."),
        ("🔑 생체 정보 및 NFC 카드 본인인증 접목", "수동 PIN 패스코드 입력 외에 학생증/사원증 NFC 하드웨어 리더 리포지토리 연동. 지문 인식 회로 등을 연계하여 높은 현장 보안성 및 원터치 편의 획득."),
        ("📈 빅데이터 기반 자원 사용 분석 대시보드", "예약 부도(No-show) 데이터 누적 분석을 통한 AI 기반 예약 패널티 및 차등 부여. 자원 배치 최적화 리포트 도출.")
    ]
    
    col_w2 = Inches(2.85)
    col_h2 = Inches(2.0)
    
    for idx, (title, desc) in enumerate(potentials):
        col_idx = idx % 2
        row_idx = idx // 2
        
        left_pos = Inches(0.8) if col_idx == 0 else Inches(3.95)
        top_pos = Inches(1.8) if row_idx == 0 else Inches(4.1)
        
        card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, col_w2, col_h2)
        card.fill.solid()
        card.fill.fore_color.rgb = PANEL_COLOR
        card.line.color.rgb = COLOR_EMERALD if idx%2==0 else COLOR_AMBER
        card.line.width = Pt(1.5)
        
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.15)
        tf_card.margin_right = Inches(0.15)
        tf_card.margin_top = Inches(0.15)
        
        p = tf_card.paragraphs[0]
        p.text = title
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_EMERALD if idx%2==0 else COLOR_AMBER
        p.space_after = Pt(4)
        
        p_desc = tf_card.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(8.5)
        p_desc.font.color.rgb = TEXT_WHITE
        p_desc.space_after = Pt(2)
        
    # AI Illustration on the Right
    try_add_image(slide10, "../images/future_expansion_illustration.png", Inches(7.1), Inches(1.8), Inches(5.4), Inches(4.3))

    # ----------------------------------------------------
    # SLIDE 11: Conclusions & Q&A
    # ----------------------------------------------------
    slide11 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide11)
    add_slide_header(slide11, "9. 결론 및 기대효과")
    
    concl_panel = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(6.8), Inches(4.5))
    concl_panel.fill.solid()
    concl_panel.fill.fore_color.rgb = PANEL_COLOR
    concl_panel.line.color.rgb = COLOR_EMERALD
    concl_panel.line.width = Pt(1.5)
    
    tf_con = concl_panel.text_frame
    tf_con.word_wrap = True
    tf_con.margin_left = Inches(0.25)
    tf_con.margin_top = Inches(0.25)
    
    p = tf_con.paragraphs[0]
    p.text = "🌟 공유 자원의 효율적 운영과 질서 있는 공동체 문화 확립"
    p.font.name = FONT_TITLE
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_EMERALD
    p.font.bold = True
    p.space_after = Pt(12)
    
    conclusions = [
        "무단 점유 및 오용 방지: 4자리 PIN 기반 자율 인증 체계로 무단 상태 해제 및 점유 가로채기 갈등 해결",
        "자원 회전 효율의 공정 분배: 1초 주기 타이머 스케줄러가 노쇼 기기를 감지해 즉시 회수하고 사용 한도 강제 해제",
        "공용 물품 관리의 범위 대중화: 회의실 뿐만 아니라 PC, 학과 공유기, 소형 키트류 등 미세 공용물도 대시보드 연동",
        "감성적 피드백을 통한 사용 의식 개선: 시인성 높은 삼색 LED 신호등 및 기분 좋은 부저 음계가 규칙 준수 유도"
    ]
    for con_text in conclusions:
        p = tf_con.add_paragraph()
        p.text = "✔ " + con_text
        p.font.name = FONT_BODY
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(4)
        p.space_after = Pt(6)
        
    # AI Illustration on the Right
    try_add_image(slide11, "../images/conclusion_illustration.png", Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.5))

    # Save
    try:
        prs.save(filename)
        print(f"Presentation saved successfully as '{filename}'.")
    except PermissionError:
        base, ext = os.path.splitext(filename)
        fallback = f"{base}_new{ext}"
        prs.save(fallback)
        print(f"Permission denied for '{filename}' (likely open in PowerPoint). Saved as '{fallback}' instead.")

if __name__ == "__main__":
    out_file = "presentation.pptx"
    if len(sys.argv) > 1:
        out_file = sys.argv[1]
    create_presentation(out_file)
